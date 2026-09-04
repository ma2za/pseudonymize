from io import BytesIO
from pathlib import Path
from typing import Any

from pseudonymize.document import ContentBlock, Document, StructuralLocation
from pseudonymize.exceptions import AdapterExecutionError

try:
    import docx
    import openpyxl
    import pptx

    HAS_OFFICE = True
except ImportError:  # pragma: no cover
    HAS_OFFICE = False


class OfficeInspectionAdapter:
    """Extracts and format-preserves Office documents (DOCX, XLSX, PPTX)."""

    def __init__(self, format_hint: str) -> None:
        if not HAS_OFFICE:
            raise RuntimeError(
                "The Office adapter requires the 'office' extra. "
                "Install it with `pip install pseudonymize[office]`."
            )
        self.format_hint = format_hint.lower()
        if self.format_hint not in ("docx", "xlsx", "pptx"):
            raise ValueError(f"Unsupported Office format hint: {format_hint}")
        self._doc: Any = None

    def extract(self, source: Path) -> Document:
        blocks: list[ContentBlock] = []
        try:
            if self.format_hint == "docx":
                self._extract_docx(source, blocks)
            elif self.format_hint == "xlsx":
                self._extract_xlsx(source, blocks)
            elif self.format_hint == "pptx":
                self._extract_pptx(source, blocks)
        except Exception:
            raise AdapterExecutionError(
                f"input adapter failed while reading the {self.format_hint.upper()} document"
            ) from None

        metadata = {"format": self.format_hint}
        return Document("file", tuple(blocks), metadata)

    def _extract_core_properties(self, doc_obj: Any, blocks: list[ContentBlock]) -> None:
        props = getattr(doc_obj, "core_properties", getattr(doc_obj, "properties", None))
        if not props:
            return
        for attr in ["author", "title", "subject", "comments", "last_modified_by", "creator"]:
            val = getattr(props, attr, None)
            if val and isinstance(val, str):
                blocks.append(
                    ContentBlock(
                        id=f"core-prop-{attr}",
                        text=val.strip(),
                        location=StructuralLocation(path=("core_properties", attr)),
                    )
                )

    def _extract_docx(self, source: Path, blocks: list[ContentBlock]) -> None:
        self._doc = docx.Document(str(source))
        self._walk_docx(self._doc, (), blocks)
        for s_idx, section in enumerate(self._doc.sections):
            self._walk_docx(section.header, ("header", s_idx), blocks)
            self._walk_docx(section.footer, ("footer", s_idx), blocks)
        self._extract_core_properties(self._doc, blocks)

    def _walk_docx(
        self, container: Any, path_prefix: tuple[Any, ...], blocks: list[ContentBlock]
    ) -> None:
        for para_index, para in enumerate(container.paragraphs):
            text = para.text.strip()
            if text:
                path = (*path_prefix, "paragraph", para_index)
                blocks.append(
                    ContentBlock(
                        id="docx-" + "-".join(str(p) for p in path),
                        text=text,
                        location=StructuralLocation(path),
                    )
                )
        for table_index, table in enumerate(container.tables):
            for row_index, row in enumerate(table.rows):
                for col_index, cell in enumerate(row.cells):
                    cell_path = (
                        *path_prefix,
                        "table",
                        table_index,
                        "row",
                        row_index,
                        "col",
                        col_index,
                    )
                    # A cell itself can contain text directly (in its paragraphs)
                    # or more nested tables. We recurse into the cell.
                    self._walk_docx(cell, cell_path, blocks)

    def _extract_xlsx(self, source: Path, blocks: list[ContentBlock]) -> None:
        self._doc = openpyxl.load_workbook(filename=source, data_only=True)
        self._extract_core_properties(self._doc, blocks)
        for sheet_index, sheet_name in enumerate(self._doc.sheetnames):
            sheet = self._doc[sheet_name]
            for row_index, row in enumerate(sheet.iter_rows(values_only=True)):
                for col_index, cell_value in enumerate(row):
                    if cell_value is not None:
                        text = str(cell_value).strip()
                        if text:
                            blocks.append(
                                ContentBlock(
                                    id=f"sheet-{sheet_index:06d}-row-{row_index:06d}-col-{col_index:06d}",
                                    text=text,
                                    location=StructuralLocation(
                                        ("sheet", sheet_name, "row", row_index, "column", col_index)
                                    ),
                                )
                            )

    def _extract_pptx(self, source: Path, blocks: list[ContentBlock]) -> None:
        self._doc = pptx.Presentation(str(source))
        self._extract_core_properties(self._doc, blocks)
        for slide_index, slide in enumerate(self._doc.slides):
            for shape_index, shape in enumerate(slide.shapes):
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        blocks.append(
                            ContentBlock(
                                id=f"slide-{slide_index:06d}-shape-{shape_index:06d}",
                                text=text,
                                location=StructuralLocation(
                                    ("slide", slide_index, "shape", shape_index)
                                ),
                            )
                        )

    def render(self, document: Document) -> bytes:
        if not self._doc:
            raise AdapterExecutionError("Cannot render before extraction")

        try:
            if self.format_hint == "docx":
                self._render_docx(document)
            elif self.format_hint == "xlsx":
                self._render_xlsx(document)
            elif self.format_hint == "pptx":
                self._render_pptx(document)
        except Exception:
            raise AdapterExecutionError(
                f"output adapter failed while rendering the {self.format_hint.upper()} document"
            ) from None

        out = BytesIO()
        self._doc.save(out)
        return out.getvalue()

    def _apply_core_properties(self, doc_obj: Any, document: Document) -> None:
        """Apply modified metadata and sanitize missing/unchanged fields."""
        props = getattr(doc_obj, "core_properties", getattr(doc_obj, "properties", None))
        if not props:
            return

        updates = {}
        for block in document.blocks:
            loc = block.location
            if (
                isinstance(loc, StructuralLocation)
                and loc.path
                and loc.path[0] == "core_properties"
            ):
                updates[str(loc.path[1])] = block.text

        for attr in ["author", "title", "subject", "comments", "last_modified_by", "creator"]:
            if not hasattr(props, attr):
                continue
            if attr in updates:
                setattr(props, attr, updates[attr])
            else:
                default = (
                    "Sanitized Document"
                    if attr == "title"
                    else "pseudonymize"
                    if attr in ("author", "creator", "last_modified_by")
                    else ""
                )
                setattr(props, attr, default)

    def _render_docx(self, document: Document) -> None:
        self._apply_core_properties(self._doc, document)
        for block in document.blocks:
            loc = block.location
            if not isinstance(loc, StructuralLocation):
                continue
            path = loc.path

            curr = self._doc
            i = 0
            while i < len(path):
                tag = path[i]
                if tag == "paragraph":
                    para_index = int(path[i + 1])
                    curr.paragraphs[para_index].text = block.text
                    break
                elif tag == "table":
                    table_index = int(path[i + 1])
                    row_index = int(path[i + 3])
                    col_index = int(path[i + 5])
                    curr = curr.tables[table_index].rows[row_index].cells[col_index]
                    i += 6
                elif tag == "header":
                    sec_idx = int(path[i + 1])
                    curr = self._doc.sections[sec_idx].header
                    i += 2
                elif tag == "footer":
                    sec_idx = int(path[i + 1])
                    curr = self._doc.sections[sec_idx].footer
                    i += 2
                elif tag == "core_properties":
                    break
                else:
                    break

    def _render_xlsx(self, document: Document) -> None:
        self._apply_core_properties(self._doc, document)

        for block in document.blocks:
            loc = block.location
            if not isinstance(loc, StructuralLocation):
                continue
            path = loc.path
            if len(path) == 6 and path[0] == "sheet":
                sheet_name = str(path[1])
                row_index = int(path[3])
                col_index = int(path[5])
                # openpyxl uses 1-based indexing
                self._doc[sheet_name].cell(
                    row=row_index + 1, column=col_index + 1, value=block.text
                )

    def _render_pptx(self, document: Document) -> None:
        self._apply_core_properties(self._doc, document)
        for block in document.blocks:
            loc = block.location
            if not isinstance(loc, StructuralLocation):
                continue
            path = loc.path
            if len(path) == 4 and path[0] == "slide":
                slide_index = int(path[1])
                shape_index = int(path[3])
                self._doc.slides[slide_index].shapes[shape_index].text = block.text
