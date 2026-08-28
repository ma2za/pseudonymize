from pathlib import Path

import pytest

from pseudonymize import TransformationMode
from pseudonymize.document import CoordinateLocation
from pseudonymize.engine import Pseudonymizer
from pseudonymize.formats import FileFormat
from pseudonymize.inspection.pdf import PDFInspectionAdapter, _source_text_color
from pseudonymize.policy import Policy

# Optional imports for generators
try:
    import fpdf

    HAS_FPDF = True
except ImportError:
    HAS_FPDF = False

try:
    import docx
    import openpyxl
    import pptx

    HAS_OFFICE = True
except ImportError:
    HAS_OFFICE = False

pytestmark = pytest.mark.integration


@pytest.fixture
def test_pdf_path(tmp_path: Path) -> Path:
    if not HAS_FPDF:
        pytest.skip("fpdf2 not installed")
    path = tmp_path / "test.pdf"
    pdf = fpdf.FPDF()
    pdf.add_page()
    pdf.set_font("helvetica", size=12)
    pdf.cell(text="Contact me at alice@example.com for more info.")
    pdf.ln()
    pdf.cell(text="   ")
    pdf.output(str(path))
    return path


@pytest.fixture
def test_pdf_ocr_path(tmp_path: Path) -> Path:
    if not HAS_FPDF:
        pytest.skip("fpdf2 not installed")
    try:
        import PIL.Image
        import PIL.ImageDraw
        import pytesseract  # noqa: F401
    except ImportError:
        pytest.skip("pytesseract or PIL not installed")

    path = tmp_path / "test_ocr.pdf"

    # Create a simple image with text and embed it into a PDF
    img = PIL.Image.new("RGB", (400, 100), color=(255, 255, 255))
    d = PIL.ImageDraw.Draw(img)
    d.text((10, 10), "Secret email is secret@example.com", fill=(0, 0, 0))
    img_path = tmp_path / "test_image.png"
    img.save(img_path)

    pdf = fpdf.FPDF()
    pdf.add_page()
    pdf.image(str(img_path), x=10, y=10, w=100)
    pdf.output(str(path))
    return path


def test_process_pdf_ocr(test_pdf_ocr_path: Path, tmp_path: Path) -> None:
    engine = Pseudonymizer(policy=Policy.default(), mode=TransformationMode.REDACTED)
    out_path = tmp_path / "out_ocr.pdf"

    # We require tesseract installed on the system for this to work
    try:
        import pytesseract

        pytesseract.get_tesseract_version()
    except Exception:
        pytest.skip("Tesseract is not installed on the system")

    engine.process_file(test_pdf_ocr_path, out_path, format=FileFormat.PDF)
    assert out_path.exists()

    import pymupdf

    doc = pymupdf.open(out_path)
    page = doc[0]
    # Check that redaction annotation is present, or text is added
    text = page.get_text()
    assert "[REDACTED]" in text
    doc.close()


@pytest.fixture
def test_docx_path(tmp_path: Path) -> Path:
    if not HAS_OFFICE:
        pytest.skip("python-docx not installed")
    path = tmp_path / "test.docx"
    doc = docx.Document()
    doc.add_paragraph("Hello, my email is bob@example.com.")
    doc.add_paragraph("   ")  # Empty text
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Header"
    table.rows[0].cells[1].text = "Value"
    table.rows[1].cells[0].text = "Data"
    table.rows[1].cells[1].text = "   "  # Empty text
    doc.save(str(path))
    return path


@pytest.fixture
def test_xlsx_path(tmp_path: Path) -> Path:
    if not HAS_OFFICE:
        pytest.skip("openpyxl not installed")
    path = tmp_path / "test.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "User Email"
    ws["B1"] = "charlie@example.com"
    ws["C1"] = None
    ws["D1"] = "   "
    wb.save(path)
    return path


@pytest.fixture
def test_pptx_path(tmp_path: Path) -> Path:
    if not HAS_OFFICE:
        pytest.skip("python-pptx not installed")
    path = tmp_path / "test.pptx"
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Presentation"
    slide.placeholders[1].text = "Reach us at diana@example.com"
    slide.shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.RECTANGLE, 0, 0, 100, 100)  # No text
    tb = slide.shapes.add_textbox(0, 0, 100, 100)
    tb.text = "   "  # Empty text
    prs.save(str(path))
    return path


def test_inspect_pdf(test_pdf_path: Path) -> None:
    engine = Pseudonymizer(
        policy=Policy.default(),
    )
    result = engine.inspect_file(test_pdf_path, format=FileFormat.PDF)
    assert result.statistics.detections_found >= 1
    assert any(d.entity_type == "EMAIL" for d in result.detections)


def test_inspect_docx(test_docx_path: Path) -> None:
    engine = Pseudonymizer(
        policy=Policy.default(),
    )
    result = engine.inspect_file(test_docx_path, format=FileFormat.DOCX)
    assert result.statistics.detections_found >= 1
    assert any(d.entity_type == "EMAIL" for d in result.detections)


def test_inspect_xlsx(test_xlsx_path: Path) -> None:
    engine = Pseudonymizer(
        policy=Policy.default(),
    )
    result = engine.inspect_file(test_xlsx_path, format=FileFormat.XLSX)
    assert result.statistics.detections_found >= 1
    assert any(d.entity_type == "EMAIL" for d in result.detections)


def test_inspect_pptx(test_pptx_path: Path) -> None:
    engine = Pseudonymizer(
        policy=Policy.default(),
    )
    result = engine.inspect_file(test_pptx_path, format=FileFormat.PPTX)
    assert result.statistics.detections_found >= 1
    assert any(d.entity_type == "EMAIL" for d in result.detections)


def test_process_pdf(test_pdf_path: Path, tmp_path: Path) -> None:
    engine = Pseudonymizer(policy=Policy.default(), mode=TransformationMode.REDACTED)
    out_path = tmp_path / "out.pdf"
    engine.process_file(test_pdf_path, out_path, format=FileFormat.PDF)
    assert out_path.exists()

    # Verify secure redaction
    import pymupdf

    doc = pymupdf.open(out_path)
    page = doc[0]
    text = page.get_text()
    assert "[REDACTED]" in text
    assert "alice@example.com" not in text
    doc.close()


def test_process_pdf_preserves_unchanged_blocks_visually(tmp_path: Path) -> None:
    if not HAS_FPDF:
        pytest.skip("fpdf2 not installed")
    import pymupdf

    source = tmp_path / "visual-preservation.pdf"
    output = tmp_path / "visual-preservation.safe.pdf"
    generated = fpdf.FPDF()
    generated.add_page()
    generated.set_font("helvetica", size=12)
    generated.text(x=20, y=20, text="Public heading")
    generated.text(x=20, y=80, text="Contact alice@example.com")
    generated.output(str(source))

    extracted = PDFInspectionAdapter().extract(source)
    unchanged = next(block for block in extracted.blocks if block.text == "Public heading")
    assert isinstance(unchanged.location, CoordinateLocation)
    clip = pymupdf.Rect(
        unchanged.location.x0,
        unchanged.location.y0,
        unchanged.location.x1,
        unchanged.location.y1,
    )

    Pseudonymizer(mode=TransformationMode.REDACTED).process_file(source, output)

    with pymupdf.open(source) as original, pymupdf.open(output) as sanitized:
        matrix = pymupdf.Matrix(2, 2)
        original_pixels = original[0].get_pixmap(matrix=matrix, clip=clip, alpha=False).samples
        sanitized_pixels = sanitized[0].get_pixmap(matrix=matrix, clip=clip, alpha=False).samples
        assert sanitized_pixels == original_pixels
        output_text = sanitized[0].get_text()
        assert "Public heading" in output_text
        assert "alice@example.com" not in output_text
        assert "[REDACTED]" in output_text


def test_process_pdf_redacts_only_the_changed_span_inside_a_block(tmp_path: Path) -> None:
    if not HAS_FPDF:
        pytest.skip("fpdf2 not installed")
    import pymupdf

    source = tmp_path / "span-preservation.pdf"
    output = tmp_path / "span-preservation.safe.pdf"
    generated = fpdf.FPDF()
    generated.add_page()
    generated.set_font("helvetica", size=12)
    generated.text(
        x=20,
        y=40,
        text="Reference ALPHA, contact alice@example.com, status active.",
    )
    generated.output(str(source))

    with pymupdf.open(source) as original:
        prefix_clip = original[0].search_for("Reference ALPHA")[0]
        suffix_clip = original[0].search_for("status active")[0]
        matrix = pymupdf.Matrix(2, 2)
        prefix_pixels = original[0].get_pixmap(matrix=matrix, clip=prefix_clip, alpha=False).samples
        suffix_pixels = original[0].get_pixmap(matrix=matrix, clip=suffix_clip, alpha=False).samples

    Pseudonymizer(mode=TransformationMode.REDACTED).process_file(source, output)

    with pymupdf.open(output) as sanitized:
        assert (
            sanitized[0].get_pixmap(matrix=matrix, clip=prefix_clip, alpha=False).samples
            == prefix_pixels
        )
        assert (
            sanitized[0].get_pixmap(matrix=matrix, clip=suffix_clip, alpha=False).samples
            == suffix_pixels
        )
        output_text = sanitized[0].get_text()
        assert "Reference ALPHA" in output_text
        assert "status active" in output_text
        assert "alice@example.com" not in output_text
        assert "[REDACTED]" in output_text


def test_process_pdf_preserves_the_background_behind_a_redaction(tmp_path: Path) -> None:
    import pymupdf

    source = tmp_path / "colored-background.pdf"
    output = tmp_path / "colored-background.safe.pdf"
    generated = pymupdf.open()
    page = generated.new_page(width=500, height=140)
    background = (0.6, 0.6, 0.6)
    page.draw_rect(
        pymupdf.Rect(30, 35, 470, 95),
        color=background,
        fill=background,
    )
    page.insert_text(
        pymupdf.Point(45, 70),
        "Contact alice@example.com status active",
        fontsize=12,
        color=(1, 1, 1),
    )
    generated.save(source)
    generated.close()

    with pymupdf.open(source) as original:
        changed_clip = original[0].search_for("alice@example.com")[0]

    Pseudonymizer(mode=TransformationMode.REDACTED).process_file(source, output)

    with pymupdf.open(output) as sanitized:
        changed_pixels = sanitized[0].get_pixmap(clip=changed_clip, alpha=False)
        channels = changed_pixels.samples
        gray_pixels = sum(
            140 <= channels[index] <= 170
            and 140 <= channels[index + 1] <= 170
            and 140 <= channels[index + 2] <= 170
            for index in range(0, len(channels), changed_pixels.n)
        )
        assert gray_pixels > changed_pixels.width * changed_pixels.height / 2
        output_text = sanitized[0].get_text()
        assert "alice@example.com" not in output_text
        assert "[REDACTED]" in output_text
        replacement_spans = [
            span
            for block in sanitized[0].get_text("dict")["blocks"]
            for line in block.get("lines", ())
            for span in line.get("spans", ())
            if "[REDACTED]" in span["text"]
        ]
        assert replacement_spans[0]["color"] == 0xFFFFFF


def test_pdf_replacement_color_falls_back_when_spans_have_no_color() -> None:
    class PageWithoutColoredText:
        def get_text(self, _mode: str, *, clip: object) -> dict[str, object]:
            return {
                "blocks": [
                    {},
                    {
                        "lines": [
                            {},
                            {"spans": [{"color": None}]},
                        ]
                    },
                ]
            }

    assert _source_text_color(PageWithoutColoredText(), object()) == (0, 0, 0)


def test_process_docx(test_docx_path: Path, tmp_path: Path) -> None:
    engine = Pseudonymizer(policy=Policy.default(), mode=TransformationMode.REDACTED)
    out_path = tmp_path / "out.docx"
    engine.process_file(test_docx_path, out_path, format=FileFormat.DOCX)
    assert out_path.exists()

    doc = docx.Document(str(out_path))
    assert doc.paragraphs[0].text == "Hello, my email is [REDACTED]."
    assert doc.core_properties.author == "pseudonymize"
    assert doc.core_properties.last_modified_by == "pseudonymize"


def test_process_xlsx(test_xlsx_path: Path, tmp_path: Path) -> None:
    engine = Pseudonymizer(policy=Policy.default(), mode=TransformationMode.REDACTED)
    out_path = tmp_path / "out.xlsx"
    engine.process_file(test_xlsx_path, out_path, format=FileFormat.XLSX)
    assert out_path.exists()

    wb = openpyxl.load_workbook(filename=out_path, data_only=True)
    ws = wb.active
    assert ws["B1"].value == "[REDACTED]"
    assert wb.properties.creator == "pseudonymize"


def test_process_pptx(test_pptx_path: Path, tmp_path: Path) -> None:
    engine = Pseudonymizer(policy=Policy.default(), mode=TransformationMode.REDACTED)
    out_path = tmp_path / "out.pptx"
    engine.process_file(test_pptx_path, out_path, format=FileFormat.PPTX)
    assert out_path.exists()

    prs = pptx.Presentation(str(out_path))
    slide = prs.slides[0]
    assert slide.placeholders[1].text == "Reach us at [REDACTED]"
    assert prs.core_properties.author == "pseudonymize"
    assert prs.core_properties.last_modified_by == "pseudonymize"
